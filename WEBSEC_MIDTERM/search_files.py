#!/usr/bin/env python3
"""
Document Search Interface
Search through all files in a directory for keywords
"""

from pathlib import Path
import re
from flask import Flask, render_template, request, jsonify
import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

app = Flask(__name__)

# Directory to search - automatically uses the directory where this script is located
SCRIPT_DIR = Path(__file__).parent.resolve()
SEARCH_DIR = SCRIPT_DIR

def search_pdf(file_path, query):
    """Search through PDF file"""
    matches = []
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text:
                    # Find all occurrences with context
                    lines = text.split('\n')
                    for line_num, line in enumerate(lines):
                        if query.lower() in line.lower():
                            # Get context (line before and after)
                            context_start = max(0, line_num - 1)
                            context_end = min(len(lines), line_num + 2)
                            context = ' '.join(lines[context_start:context_end])
                            
                            matches.append({
                                'page': page_num,
                                'text': line.strip(),
                                'context': context[:300]  # Limit context length
                            })
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
    
    return matches

def search_docx(file_path, query):
    """Search through Word document"""
    matches = []
    try:
        doc = Document(file_path)
        for para_num, paragraph in enumerate(doc.paragraphs, 1):
            if query.lower() in paragraph.text.lower():
                matches.append({
                    'paragraph': para_num,
                    'text': paragraph.text.strip(),
                    'context': paragraph.text[:300]
                })
    except Exception as e:
        print(f"Error reading DOCX {file_path}: {e}")
    
    return matches

def search_pptx(file_path, query):
    """Search through PowerPoint presentation"""
    matches = []
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            full_text = ' '.join(slide_text)
            if query.lower() in full_text.lower():
                matches.append({
                    'slide': slide_num,
                    'text': full_text[:200],
                    'context': full_text[:300]
                })
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
    
    return matches

def search_png(file_path, query):
    """Search through PNG using OCR"""
    matches = []
    try:
        # Note: This requires tesseract-ocr to be installed
        # brew install tesseract
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        
        if query.lower() in text.lower():
            lines = text.split('\n')
            for line in lines:
                if query.lower() in line.lower():
                    matches.append({
                        'text': line.strip(),
                        'context': text[:300]
                    })
    except Exception as e:
        print(f"Error reading PNG {file_path}: {e}")
    
    return matches

def search_all_files(query):
    """Search through all supported files in directory"""
    results = []
    
    supported_extensions = {'.pdf', '.docx', '.pptx', '.png'}
    files = [f for f in SEARCH_DIR.iterdir() 
             if f.is_file() and f.suffix.lower() in supported_extensions]
    
    for file_path in files:
        ext = file_path.suffix.lower()
        matches = []
        
        if ext == '.pdf':
            matches = search_pdf(file_path, query)
        elif ext == '.docx':
            matches = search_docx(file_path, query)
        elif ext == '.pptx':
            matches = search_pptx(file_path, query)
        elif ext == '.png':
            matches = search_png(file_path, query)
        
        if matches:
            results.append({
                'filename': file_path.name,
                'filepath': str(file_path),
                'type': ext[1:].upper(),
                'matches': matches,
                'match_count': len(matches)
            })
    
    # Sort by number of matches (most matches first)
    results.sort(key=lambda x: x['match_count'], reverse=True)
    
    return results

@app.route('/')
def index():
    """Render search interface"""
    return render_template('index.html')

@app.route('/search', methods=['POST'])
def search():
    """Handle search request"""
    query = request.json.get('query', '').strip()
    
    if not query:
        return jsonify({'error': 'Please enter a search query'}), 400
    
    results = search_all_files(query)
    
    return jsonify({
        'query': query,
        'results': results,
        'total_files': len(results),
        'total_matches': sum(r['match_count'] for r in results)
    })

@app.route('/open-file', methods=['POST'])
def open_file():
    """Open a file in the default application"""
    import subprocess
    import platform
    import os
    
    filepath = request.json.get('filepath', '').strip()
    
    print(f"\n=== Open File Request ===")
    print(f"Received filepath: {filepath}")
    print(f"Operating System: {platform.system()}")
    
    if not filepath:
        return jsonify({'error': 'No file path provided'}), 400
    
    # Normalize path separators for the current OS
    # Path() handles this automatically, but we'll be explicit
    file_path = Path(filepath)
    
    print(f"Converted to Path: {file_path}")
    print(f"Path exists: {file_path.exists()}")
    print(f"Path is absolute: {file_path.is_absolute()}")
    
    # If path doesn't exist, try treating it as relative to search directory
    if not file_path.exists():
        # Try relative to search directory
        alt_path = SEARCH_DIR / file_path.name
        print(f"Trying alternate path: {alt_path}")
        print(f"Alternate path exists: {alt_path.exists()}")
        if alt_path.exists():
            file_path = alt_path
    
    if not file_path.exists():
        print(f"ERROR: File not found at any location")
        print(f"Search directory contents: {list(SEARCH_DIR.glob('*.pdf'))[:5]}...")  # Show first 5 PDFs
        return jsonify({'error': f'File not found: {file_path.name}'}), 404
    
    print(f"Opening file: {file_path}")
    print(f"File absolute path: {file_path.absolute()}")
    
    try:
        system = platform.system()
        
        if system == 'Darwin':  # macOS
            subprocess.run(['open', str(file_path)], check=True)
        elif system == 'Windows':  # Windows
            # Use the absolute path and convert to string
            os.startfile(str(file_path.absolute()))
        elif system == 'Linux':  # Linux
            subprocess.run(['xdg-open', str(file_path)], check=True)
        else:
            return jsonify({'error': f'Unsupported operating system: {system}'}), 500
            
        print(f"SUCCESS: Opened {file_path.name}")
        return jsonify({'success': True, 'message': f'Opened {file_path.name}'})
    except Exception as e:
        print(f"ERROR: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'Failed to open file: {str(e)}'}), 500

if __name__ == '__main__':
    # Check if templates directory exists
    templates_dir = Path(__file__).parent / 'templates'
    templates_dir.mkdir(exist_ok=True)
    
    print("\n" + "="*60)
    print("Document Search Interface")
    print("="*60)
    print(f"Searching in: {SEARCH_DIR}")
    print("\nStarting server at: http://127.0.0.1:5000")
    print("Press Ctrl+C to stop")
    print("="*60 + "\n")
    
    app.run(debug=True, port=5000)
